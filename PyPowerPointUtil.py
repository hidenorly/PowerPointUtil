#   Copyright 2023 hidenorly
#
#   Licensed under the Apache License, Version 2.0 (the "License");
#   you may not use this file except in compliance with the License.
#   You may obtain a copy of the License at
#
#       http://www.apache.org/licenses/LICENSE-2.0
#
#   Unless required by applicable law or agreed to in writing, software
#   distributed under the License is distributed on an "AS IS" BASIS,
#   WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#   See the License for the specific language governing permissions and
#   limitations under the License.

import os
import string
import time

import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor
import webcolors

from io import BytesIO


class PowerPointUtil:
    SLIDE_WIDTH_INCH = 16
    SLIDE_HEIGHT_INCH = 9

    def __init__(self, path):
        self.prs = Presentation()
        self.prs.slide_width  = Inches(self.SLIDE_WIDTH_INCH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT_INCH)
        self.path = path
        self.currentSlide = None

    def save(self):
        self.prs.save(self.path)

    # layout is full, left, right, top, bottom
    def getLayoutPosition(self, layout="full"):
        # for full
        x=0
        y=0
        width = self.prs.slide_width
        height = self.prs.slide_height

        if layout=="left" or layout=="right":
            width = width /2
        if layout=="top" or layout=="bottom":
            height = height /2
        if layout=="right":
            x=width
        if layout=="bottom":
            y=height

        return x,y,width,height

    def getLayoutToFitRegion(self, width, height, regionWidth, regionHeight):
        resultWidth = width
        resultHeight = height

        if width > height:
            resultWidth = regionWidth
            resultHeight = int(regionWidth * height / width+0.99)
        else:
            resultHeight = regionHeight
            resultWidth = int(regionHeight * width / height+0.99)

        return resultWidth,regionHeight

    def getLayoutWithinRegion(self, width, height, regionWidth, regionHeight, isFitWihthinRegion=True):
        ratio = min(regionWidth/self.prs.slide_width, regionHeight/self.prs.slide_height)
        resultWidth = float(width) * float(ratio)
        resultHeight = float(height) * float(ratio)

        if isFitWihthinRegion:
            deltaWidth = resultWidth - regionWidth
            deltaHeight = resultHeight - regionHeight
            if deltaWidth>0 or deltaHeight>0:
                # exceed the region then scale is required
                if deltaWidth>deltaHeight:
                    resultWidth = regionWidth
                    ratio =  regionWidth/resultWidth
                    resultHeight = resultHeight * ratio + 0.99
                else:
                    resultHeight = regionHeight
                    ratio = regionHeight / resultHeight
                    resultWidth = resultWidth * ratio

        return int(resultWidth+0.99), int(resultHeight+0.99)


    def addSlide(self, layout=None):
        if layout == None:
            layout = self.prs.slide_layouts[6]
        self.currentSlide = self.prs.slides.add_slide(layout)

        return self.currentSlide

    def copySlideContent(self, srcSlide, dstSlide=None, layout=None, enableLayoutOverride=False):
        regionX, regionY, regionWidth, regionHeight = self.getLayoutPosition(layout)

        if dstSlide == None:
            dstSlide = self.currentSlide
        if dstSlide:
            if isinstance(srcSlide, PowerPointUtil):
                srcSlide = srcSlide.currentSlide
            for srcShape in srcSlide.shapes:
                nextX = srcShape.left + regionX
                nextY = srcShape.top + regionY
                # TODO: ensure nextX/Y in the region
                nextWidth, nextHeight = self.getLayoutWithinRegion(srcShape.width, srcShape.height, regionWidth, regionHeight)

                if srcShape.has_text_frame:
                    newShape = dstSlide.shapes.add_textbox(
                        nextX, nextY, nextWidth, nextHeight
                    )
                    PowerPointUtil.copyTextFormat(srcShape.text_frame, newShape.text_frame)

                    if enableLayoutOverride:
                        alignment = PP_ALIGN.LEFT
                        if layout == "right":
                            alignment = PP_ALIGN.RIGHT
                        for paragraph in newShape.text_frame.paragraphs:
                            paragraph.alignment = alignment

                elif srcShape.shape_type == 13:  # Shape type 13 : Picture
                    img = srcShape.image
                    imgData = img.blob
                    imgStream = BytesIO(imgData)
                    newShape = dstSlide.shapes.add_picture(imgStream, nextX, nextY, nextWidth, nextHeight)
                    newShape.width = int(nextWidth+0.99)
                    newShape.height = int(nextHeight+0.99)


    # --- picture
    def addPicture(self, imagePath, x=0, y=0, width=None, height=None, isFitToSlide=True, regionWidth=None, regionHeight=None, isFitWihthinRegion=False):
        if not regionWidth:
            regionWidth = self.prs.slide_width
        if not regionHeight:
            regionHeight = self.prs.slide_height
        regionWidth = int(regionWidth+0.99)
        regionHeight = int(regionHeight+0.99)
        pic = None
        try:
            if self.currentSlide:
                pic = self.currentSlide.shapes.add_picture(imagePath, x, y)
        except:
            pass
        if pic:
            if width and height:
                pic.width = width
                pic.height = height
            else:
                if isFitToSlide:
                    width, height = pic.image.size
                    picWidth = pic.width
                    picHeight = pic.height
                    if width > height:
                        picWidth = regionWidth
                        picHeight = int(regionWidth * height / width + 0.99)
                    else:
                        picHeight = regionHeight
                        picWidth = int(regionHeight * width / height + 0.99)
                    if isFitWihthinRegion:
                        deltaWidth = picWidth - regionWidth
                        deltaHeight = picHeight - regionHeight
                        if deltaWidth>0 or deltaHeight>0:
                            # exceed the region
                            if deltaWidth > deltaHeight:
                                picWidth = regionWidth
                                picHeight = int(regionWidth * height / width + 0.99)
                            else:
                                picHeight = regionHeight
                                picWidth = int(regionHeight * width / height + 0.99)
                    pic.width = picWidth
                    pic.height = picHeight
        return pic

    # --- text ----
    @staticmethod
    def nameToRgb(name):
        result = RGBColor(0,0,0)
        try:
            rgb = webcolors.name_to_rgb(name)
            result = RGBColor(rgb.red, rgb.green, rgb.blue)
        except:
            pass
        return result

    @staticmethod
    def applyExFormat(exFormat, textbox, font, text_frame):
        exFormats = exFormat.split(",")
        for anFormat in exFormats:
            cmdarg = anFormat.split(":")
            cmd = cmdarg[0]
            val = None
            if len(cmdarg)>=2:
                val = cmdarg[1]
            if cmd=="color":
                font.color.rgb = PowerPointUtil.nameToRgb(val)
            elif cmd=="face":
                font.name = val
            elif cmd=="size":
                font.size = Pt(float(val))
            elif cmd=="bold":
                font.bold = True
            elif cmd=="effect":
                # TODO: fix
                shadow = textbox.shadow
                shadow.visible = True
                shadow.shadow_type = 'outer'
                shadow.style = 'outer'
                shadow.blur_radius = Pt(5)
                shadow.distance = Pt(2)
                shadow.angle = 45
                shadow.color = MSO_THEME_COLOR_INDEX.ACCENT_5
                shadow.transparency = 0

    @staticmethod
    def copyTextFormat(srcTextFrame, dstTextFrame, ):
        i = 0
        dstParagraph = dstTextFrame.paragraphs[0]
        for srcParagraph in srcTextFrame.paragraphs:
            if i!=0:
                dstParagraph = dstTextFrame.add_paragraph()

            if srcParagraph.font.size is not None:
                dstParagraph.font.size = srcParagraph.font.size
            if srcParagraph.font.name is not None:
                dstParagraph.font.name = srcParagraph.font.name
            dstParagraph.text = srcParagraph.text
            if srcParagraph.alignment is not None:
                dstParagraph.alignment = srcParagraph.alignment
            if srcParagraph.font.color is not None and hasattr(srcParagraph.font.color, 'rgb') and srcParagraph.font.color.rgb is not None:
                dstParagraph.font.color.rgb = srcParagraph.font.color.rgb
            if srcParagraph.font.bold is not None:
                dstParagraph.font.bold = srcParagraph.font.bold
            if srcParagraph.font.italic is not None:
                dstParagraph.font.italic = srcParagraph.font.italic
            if srcParagraph.font.underline is not None:
                dstParagraph.font.underline = srcParagraph.font.underline
            i = i + 1

    def addText(self, text, x=Inches(0), y=Inches(0), width=None, height=None, fontFace='Calibri', fontSize=Pt(18), isAdjustSize=True, textAlign = PP_ALIGN.LEFT, isVerticalCenter=False, exFormat=None):
        if width==None:
            width=self.prs.slide_width
        if height==None:
            height=self.prs.slide_height
        width = int(width+0.99)
        height = int(height+0.99)

        if self.currentSlide:
            textbox = self.currentSlide.shapes.add_textbox(x, y, width, height)
            text_frame = textbox.text_frame
            text_frame.text = text
            font = text_frame.paragraphs[0].font
            font.name = fontFace
            font.size = fontSize
            theHeight = textbox.height

            if exFormat:
                PowerPointUtil.applyExFormat(exFormat, textbox, font, text_frame)

            if isAdjustSize:
                text_frame.auto_size = True
                textbox.top = y

            if isVerticalCenter:
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = textAlign
