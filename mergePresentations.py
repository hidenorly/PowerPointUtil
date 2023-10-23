#   Copyright 2023 hidenorly
#
#   Licensed under the Apache License, Version 2.0 (the "License");
#   you may not use this file except in compliance with the License.
#   You may obtain a copy of the License at
#
#       http://www.apache.org/licenses/LICENSE-2.0
#
#   Unless required by applicable law or agreed to in writing, software
#   distributed under the License is distributed on an "AS IS" BASIS,
#   WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#   See the License for the specific language governing permissions and
#   limitations under the License.

import argparse
import collections.abc
from pptx import Presentation
from PyPowerPointUtil import PowerPointUtil


def mergePowerpoints(inputFiles, outputFile, layouts=None, enableLayoutOverride=False):
    srcPresentations = [Presentation(inputFile) for inputFile in inputFiles]
    maxPages = max(len(presentation.slides) for presentation in srcPresentations)
    mergedPresentation = PowerPointUtil(outputFile)
    newSlides = []
    layouts = str(layouts).split(",")
    if len(srcPresentations)>len(layouts):
        layouts.extend([""] * ( len(srcPresentations) - len(layouts) ) )

    for page in range(maxPages):
        newSlides.append( mergedPresentation.addSlide() )

    index = 0
    for srcPresentation in srcPresentations:
        layout = layouts[index]
        i = 0        
        for srcSlide in srcPresentation.slides:
            mergedPresentation.copySlideContent(srcSlide, newSlides[i], layout, enableLayoutOverride)
            i = i + 1
        index = index + 1
    
    mergedPresentation.save()


def concatPowerpoints(inputFiles, outputFile):
    srcPresentations = [Presentation(inputFile) for inputFile in inputFiles]
    mergedPresentation = PowerPointUtil(outputFile)
    newSlides = []

    for srcPresentation in srcPresentations:
        for srcSlide in srcPresentation.slides:
            mergedPresentation.addSlide()
            mergedPresentation.copySlideContent(srcSlide)
    
    mergedPresentation.save()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Merge PowerPoint files. Usage: -i ppt1.pptx -i ppt2.pptx -o output.pptx")
    parser.add_argument("-i", "--input", required=True, action='append', default=[], help="Input PowerPoint files")
    parser.add_argument("-o", "--output", required=True, help="Output PowerPoint file")
    parser.add_argument("-m", "--mode", default="combine", help="set combine or append")
    parser.add_argument("-l", "--layouts", default="", help="set left or right or top or bottom or \"\"")
    parser.add_argument("-r", "--override", default=False, action='store_true', help="set if override the layout as the -l specified layout")
    args = parser.parse_args()

    if args.mode == "combine":
        mergePowerpoints(args.input, args.output, args.layouts, args.override)
    else:
        concatPowerpoints(args.input, args.output)
